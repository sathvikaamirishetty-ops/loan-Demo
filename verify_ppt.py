from pptx import Presentation
import os

try:
    file_path = 'Loan_Eligibility_Demo.pptx'
    if not os.path.exists(file_path):
        print("File does not exist.")
    else:
        size = os.path.getsize(file_path)
        print(f"File size: {size} bytes")
        
        prs = Presentation(file_path)
        print(f"Successfully loaded. Slide count: {len(prs.slides)}")
        for i, slide in enumerate(prs.slides):
            try:
                title = slide.shapes.title.text
            except:
                title = "No Title"
            print(f"Slide {i+1}: {title}")

except Exception as e:
    print(f"Error loading presentation: {e}")
