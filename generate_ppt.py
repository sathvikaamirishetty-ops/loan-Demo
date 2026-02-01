from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()

    # Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Loan Eligibility Comparison Demo"
    subtitle.text = "A Side-by-Side Bank Eligibility Estimator\nProposed Concept"

    # Slide 2: The Problem
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Current Market Gap"
    tf = body_shape.text_frame
    tf.text = "Most existing tools focus on:"
    p = tf.add_paragraph()
    p.text = "Interest Rate comparisons"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "EMI calculations for a fixed loan amount"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Lead generation rather than instant transparency"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\nWhat is missing?"
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = "A simple 'How much can I get?' comparison across banks without entering a specific loan amount request."
    p.level = 1

    # Slide 3: The Solution
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Our Solution: 2-Bank Dummy Demo"
    tf = body_shape.text_frame
    tf.text = "Key Features:"
    p = tf.add_paragraph()
    p.text = "Single Input: User enters their financial profile once."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Direct Comparison: Side-by-side view of max loan eligibility."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Dummy Logic: Uses indicative FOIR and Multiplier rules to simulate realistic differences."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Immediate Value: Users see 'HDFC gives ₹58L vs BoB gives ₹38L' instantly."
    p.level = 1

    # Slide 4: Logic & Rules (HDFC vs BoB)
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Behind the Scenes: The Logic"
    tf = body_shape.text_frame
    tf.text = "We use simplified banking norms for the demo:"
    
    p = tf.add_paragraph()
    p.text = "\nHDFC Bank (Aggressive Profile)"
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = "FOIR: 55% | Multiplier: 200x"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Result: Higher loan amount eligibility."
    p.level = 1

    p = tf.add_paragraph()
    p.text = "\nBank of Baroda (Conservative Profile)"
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = "FOIR: 45% | Multiplier: 180x"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Result: Lower loan amount, stricter criteria."
    p.level = 1

    # Slide 5: Tech Stack
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Technology Stack"
    tf = body_shape.text_frame
    tf.text = "Simple, lightweight, and effective:"
    p = tf.add_paragraph()
    p.text = "Backend: Python (Flask)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Frontend: HTML5 + Vanilla CSS (Premium Styling)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Deployment: Ready for local demo or cloud hosting"
    p.level = 1

    # Slide 6: Demo Walkthrough
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Live Demo Flow"
    tf = body_shape.text_frame
    tf.text = "1. Enter Income & EMIs"
    p = tf.add_paragraph()
    p.text = "2. Select Employment & Credit Score"
    p = tf.add_paragraph()
    p.text = "3. Click 'Check Eligibility'"
    p = tf.add_paragraph()
    p.text = "4. View Comparative Results"
    
    p = tf.add_paragraph()
    p.text = "\nStatus: Verified & Running"
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 128, 0)

    prs.save('Loan_Eligibility_Demo.pptx')
    print("Presentation saved successfully.")

if __name__ == "__main__":
    create_presentation()
